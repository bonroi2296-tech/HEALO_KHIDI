"""
리플렛 PPTX 조립 v2 — 사진·색면·아이콘·글자를 «전부 개별 객체»로

  python build-leaflet-pptx.py <작업폴더> <접두어> <출력.pptx>

파워포인트에서 사진은 「그림 바꾸기」로, 색면은 「도형 채우기」로, 글자는 클릭해서 고칠 수 있다.
"""
import json, sys, os
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

WORK, TAG, OUT = sys.argv[1], sys.argv[2], sys.argv[3]
PARTS = os.path.join(WORK, f"{TAG}-parts")
SHEET_W_MM, SHEET_H_MM = 303.0, 216.0
NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Pretendard 는 굵기마다 «패밀리 이름»이 다르게 설치돼 있다(실측).
# 700 을 「Pretendard + 굵게」로 주면 파워포인트가 Bold 페이스를 못 써서 얇게 나온다.
FONT_BY_WEIGHT = {
    300: ("Pretendard Light", False),
    400: ("Pretendard", False),
    500: ("Pretendard Medium", False),
    600: ("Pretendard SemiBold", False),
    700: ("Pretendard SemiBold", True),
    800: ("Pretendard ExtraBold", False),
    900: ("Pretendard ExtraBold", True),
}
# 명조(제목용) — 굵기마다 패밀리 이름이 따로 설치돼 있다
SERIF_BY_WEIGHT = {
    300: ("KoPubWorldBatang_Pro Light", False),
    400: ("KoPubWorldBatang_Pro Light", False),
    500: ("KoPubWorldBatang_Pro Medium", False),
    600: ("KoPubWorldBatang_Pro Medium", False),
    700: ("KoPubWorldBatang_Pro Bold", False),
    800: ("KoPubWorldBatang_Pro Bold", False),
}

def font_for(w, fam=None):
    table = SERIF_BY_WEIGHT if (fam and "Serif" in fam) else FONT_BY_WEIGHT
    k = min(table, key=lambda x: abs(x - w))
    return table[k]

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
rgb = lambda h: RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def set_alpha(color_el_parent, alpha):
    """채우기·선에 투명도를 넣는다 (파워포인트는 1/1000 % 단위)"""
    if alpha is None or alpha >= 0.999:
        return
    fill = color_el_parent.find(f"{NS}solidFill")
    if fill is None:
        return
    clr = fill.find(f"{NS}srgbClr")
    if clr is None:
        return
    a = etree.SubElement(clr, f"{NS}alpha")
    a.set("val", str(int(round(alpha * 100000))))

data = json.load(open(os.path.join(WORK, f"{TAG}-parts.json"), encoding="utf-8"))

prs = Presentation()
prs.slide_width, prs.slide_height = Mm(SHEET_W_MM), Mm(SHEET_H_MM)
blank = prs.slide_layouts[6]

def patch_theme(prs):
    for part in prs.part.package.iter_parts():
        if "theme" not in part.partname.lower():
            continue
        root = etree.fromstring(part.blob)
        changed = False
        for scheme in root.iter(f"{NS}fontScheme"):
            for tag in ("majorFont", "minorFont"):
                fn = scheme.find(f"{NS}{tag}")
                if fn is None:
                    continue
                for sub in ("latin", "ea", "cs"):
                    el = fn.find(f"{NS}{sub}")
                    if el is not None:
                        el.set("typeface", "Pretendard"); changed = True
        if changed:
            part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

def set_run_font(run, name, bold, size_pt, color, spc_pt):
    f = run.font
    f.name, f.bold, f.size = name, bold, Pt(size_pt)
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):          # 한글·복합문자 칸도 못 박는다
        el = rPr.find(f"{NS}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{NS}{tag}")
        el.set("typeface", name)
    if spc_pt:
        rPr.set("spc", str(int(round(spc_pt * 100))))

def add_shape(slide, s):
    x, y, w, h = Mm(s["xMm"]), Mm(s["yMm"]), Mm(max(s["wMm"], 0.1)), Mm(max(s["hMm"], 0.1))
    kind = s["kind"]
    if kind == "line":
        t = max(s.get("thickMm", 0.2), 0.18)
        side = s["side"]
        if side == "top":    box = (s["xMm"], s["yMm"], s["wMm"], t)
        elif side == "bottom": box = (s["xMm"], s["yMm"] + s["hMm"] - t, s["wMm"], t)
        elif side == "left":  box = (s["xMm"], s["yMm"], t, s["hMm"])
        else:                 box = (s["xMm"] + s["wMm"] - t, s["yMm"], t, s["hMm"])
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(box[0]), Mm(box[1]), Mm(box[2]), Mm(box[3]))
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(s["color"])
        set_alpha(sh.fill._xPr.find(f"{NS}solidFill").getparent(), s.get("alpha", 1))
        sh.line.fill.background(); sh.shadow.inherit = False
        sh.text_frame.word_wrap = False
        return sh

    if kind == "oval":
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    elif s.get("radiusMm", 0) > 0.25:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        try:
            sh.adjustments[0] = min(0.5, s["radiusMm"] / max(min(s["wMm"], s["hMm"]), 0.1))
        except Exception:
            pass
    else:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)

    if s.get("fill"):
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(s["fill"])
        set_alpha(sh.fill._xPr.find(f"{NS}solidFill").getparent(), s.get("fillAlpha", 1))
    else:
        sh.fill.background()
    ln = s.get("line")
    if ln:
        sh.line.color.rgb = rgb(ln["color"])
        sh.line.width = Pt(max(ln["thickMm"] * 72 / 25.4, 0.5))
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = False
    return sh

def add_text(slide, blk):
    # 세로 가운데 정렬 칸(원 안의 번호 등)은 여유를 붙이면 «가운데»가 어긋난다 — 잰 그대로 쓴다.
    mid = blk.get("flexAlign") == "center"
    px, py, pw, ph = (0, 0, 0, 0) if mid else (0.15, 0.5, 0.9, 1.2)
    tb = slide.shapes.add_textbox(Mm(blk["xMm"] - px), Mm(blk["yMm"] - py),
                                  Mm(blk["wMm"] + pw), Mm(blk["hMm"] + ph))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if blk.get("flexJustify") == "center":
        blk["align"] = "center"
    if blk.get("flexAlign") == "center":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = ALIGN.get(blk["align"], PP_ALIGN.LEFT)
    # 줄 높이를 못 박으면 파워포인트가 글자를 아래로 흘린다 — 가운데 정렬 칸에서는 글꼴 기본에 맡긴다.
    if blk["lineHeightPt"] and not mid:
        para.line_spacing = Pt(blk["lineHeightPt"])
    cur = para
    for rn in blk["runs"]:
        if rn.get("br"):
            cur = tf.add_paragraph(); cur.alignment = para.alignment
            if blk["lineHeightPt"] and not mid:
                cur.line_spacing = Pt(blk["lineHeightPt"])
            continue
        if not rn.get("text"):
            continue
        r = cur.add_run(); r.text = rn["text"]
        name, bold = font_for(rn["weight"], rn.get("fam"))
        col = rgb(rn["color"])
        set_run_font(r, name, bold, rn["sizePt"], col, rn.get("spcPt", 0))
        if rn.get("alpha", 1) < 0.999:
            rPr = r._r.get_or_add_rPr()
            set_alpha(rPr, rn["alpha"])
    return tb

counts = {"shape": 0, "image": 0, "icon": 0, "text": 0}
for si in range(data["sheetCount"]):
    slide = prs.slides.add_slide(blank)
    items = []
    for s in data["shapes"]:
        if s["sheet"] == si: items.append((s["z"], "shape", s))
    for m in data["images"]:
        if m["sheet"] == si: items.append((m["z"], "image", m))
    for ic in data["icons"]:
        if ic["sheet"] == si: items.append((ic["z"], "icon", ic))
    for t in data["texts"]:
        if t["sheet"] == si: items.append((t["z"], "text", t))
    items.sort(key=lambda x: x[0])          # 그리는 순서 = 원본 겹침 순서

    for _, kind, it in items:
        if kind == "shape":
            add_shape(slide, it); counts["shape"] += 1
        elif kind in ("image", "icon"):
            p = os.path.join(PARTS, f"{it['id']}.png")
            if not os.path.exists(p):
                print(f"  [빠짐] {it['id']} 파일 없음"); continue
            slide.shapes.add_picture(p, Mm(it["xMm"]), Mm(it["yMm"]), Mm(it["wMm"]), Mm(it["hMm"]))
            counts["image" if kind == "image" else "icon"] += 1
        else:
            add_text(slide, it); counts["text"] += 1

patch_theme(prs)
prs.save(OUT)
print(f"  도형 {counts['shape']} · 그림 {counts['image']} · 아이콘 {counts['icon']} · 글자상자 {counts['text']}")
print(f"  → {os.path.basename(OUT)}  ({round(os.path.getsize(OUT)/1024/1024, 2)} MB)")
