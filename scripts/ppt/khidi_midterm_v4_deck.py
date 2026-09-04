# -*- coding: utf-8 -*-
"""KHIDI 중간평가 발표자료 v4 — 밀도를 낮춘 판

2026-08-31 PO: "그냥 다 빼고 니 기본 스타일대로 만들어봐. 색도 다 무시해도 괜찮음"
→ BeyondK 깔(teal)과 표 중심 구성을 놓고 처음부터 다시 짠다.

무엇이 다른가
  · 한 쪽에 한 문장. 제목이 곧 그 쪽의 결론이다.
  · 표를 화면에서 뺐다. 숫자는 크게, 근거는 별첨으로 돌린다.
  · 색을 안 쓴다. 먹색·회색·흰색 세 단계와 「검은 판」만으로 강약을 만든다.
  · 28쪽 → 16쪽.

숫자: 2026-08-31 재검수판 (문의 6건 · 사전상담 29건 · 파트너 미팅 23회 · 계약 10건)
"""
import os
import sys

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

IMG = sys.argv[1] if len(sys.argv) > 1 else "그림"
OUT = sys.argv[2] if len(sys.argv) > 2 else "중간평가_발표자료_v4.pptx"

# ── 색: 세 단계뿐. 강조는 색이 아니라 「검은 판」이 한다
INK = RGBColor(0x11, 0x11, 0x11)      # 본문 먹색
GRAY = RGBColor(0x88, 0x88, 0x88)     # 보조
FAINT = RGBColor(0xCF, 0xCF, 0xCF)    # 가는 선
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2B, 0x2B, 0x2B)     # 짙은 판 바탕 (표지·마무리에만)
ONDARK = RGBColor(0xF5, 0xF5, 0xF5)   # 검은 판 위 글씨
DIM = RGBColor(0x9A, 0x9A, 0x9A)      # 검은 판 위 보조

HEAVY, XBOLD, MED, REG, LIGHT = (
    "에스코어 드림 8 Heavy", "에스코어 드림 7 ExtraBold",
    "에스코어 드림 5 Medium", "에스코어 드림 4 Regular", "에스코어 드림 3 Light")

W, H = 960, 540
M = 84                                 # 넉넉한 바깥 여백


def pt(v):
    return Emu(int(v * 12700))


prs = Presentation()
prs.slide_width, prs.slide_height = pt(W), pt(H)
BLANK = prs.slide_layouts[6]


def 장(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, pt(W), pt(H))
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK if dark else PAPER
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def 글(s, x, y, w, h, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(pt(x), pt(y), pt(w), pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.paragraphs[0].alignment = align
    return tf


def 줄(tf, txt, size, color=INK, font=REG, first=False, before=0, spacing=None, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_before = Pt(before)
    if spacing:
        p.line_spacing = spacing
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.name = font; r.font.color.rgb = color
    글꼴박기(r, font)
    return p


def 글꼴박기(run, 이름):
    """python-pptx 의 font.name 은 <a:latin> 만 건드린다.

    한글은 <a:ea>(동아시아) 를 따라가고, 그게 비면 테마의 script="Hang" 값으로 떨어진다.
    빈 발표자료의 그 값이 「맑은 고딕」이라 새로 친 글자가 대체된다(2026-08-31 실측).
    그래서 latin·ea·cs 를 함께 박는다."""
    rPr = run._r.get_or_add_rPr()
    for 태그 in ("a:ea", "a:cs"):
        el = rPr.find(qn(태그))
        if el is None:
            el = etree.SubElement(rPr, qn(태그))
        el.set("typeface", 이름)


def 선(s, x, y, w, color=FAINT, 두께=0.8):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, pt(x), pt(y), pt(w), pt(두께))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False


def 상자(s, x, y, w, h, color):
    """바탕색만 있는 사각형. 숫자 카드에 쓴다."""
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, pt(x), pt(y), pt(w), pt(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def 그림(s, 파일, x, y, w=None, h=None):
    f = os.path.join(IMG, 파일)
    if not os.path.exists(f):
        return None
    kw = {}
    if w: kw["width"] = pt(w)
    if h: kw["height"] = pt(h)
    return s.shapes.add_picture(f, pt(x), pt(y), **kw)


def 머리(s, 제목, 부제=None, dark=False):
    """제목이 곧 그 쪽의 결론이다. 그래서 크게 쓴다."""
    색 = ONDARK if dark else INK
    보조 = DIM if dark else GRAY
    tf = 글(s, M, 62, W - M * 2, 76)
    줄(tf, 제목, 25, 색, XBOLD, first=True, spacing=1.16)
    if 부제:
        줄(tf, 부제, 12, 보조, REG, before=8)
    return s


def 쪽번호(s, n, dark=False):
    tf = 글(s, W - M - 60, H - 46, 60, 20, align=PP_ALIGN.RIGHT)
    줄(tf, str(n), 9.5, DIM if dark else FAINT, MED, first=True)


def 각주(s, txt, dark=False):
    tf = 글(s, M, H - 48, W - M * 2 - 80, 22)
    줄(tf, txt, 9, DIM if dark else GRAY, LIGHT, first=True)


번호 = [0]
def 세기(s, dark=False):
    번호[0] += 1
    if 번호[0] > 1:
        쪽번호(s, 번호[0], dark)
    return s


# ══════════════════════════════════════════ 1. 표지
s = 세기(장(dark=True), dark=True)
tf = 글(s, M, 196, 700, 40)
줄(tf, "2026년 ICT기반 외국인환자 사전상담·사후관리 지원 사업", 12, DIM, REG, first=True)
tf = 글(s, M, 232, 760, 90)
줄(tf, "중간보고", 52, ONDARK, HEAVY, first=True)
선(s, M, 344, 96, DIM, 1.6)
tf = 글(s, M, 368, 720, 60)
줄(tf, "카자흐스탄 암환자 대상 ICT 기반 사전상담·사후관리 통합 의료연계 서비스", 13, ONDARK, REG, first=True)
줄(tf, "healwith 플랫폼", 13, DIM, REG, before=4)
tf = 글(s, M, H - 76, 760, 24)
줄(tf, "주관 본로이   ·   참여 면력한방병원 · 신촌면력한방병원   ·   2026. 09.", 10, DIM, REG, first=True)

# ══════════════════════════════════════════ 2. 한 장 요약
s = 세기(장())
머리(s, "추진 현황 총괄", "여덟 달 중 다섯 달이 지났습니다. 만든 것은 다 됐고, 환자를 데려오는 일이 남았습니다")
넷 = [("플랫폼", "구축 완료", "6개 언어 · 원격협진 · AI 상담"),
      ("파트너", "23곳 만나 10곳 계약", "카자흐 · 러시아 · 키르기스 · 인도"),
      ("사전상담", "29건", "환자 6명 · 목표 120건의 24%"),
      ("유치", "0건", "목표 12건")]
y = 186
for i, (칸, 값, 설명) in enumerate(넷):
    if i:
        선(s, M, y - 22, W - M * 2)
    tf = 글(s, M, y, 130, 40)
    줄(tf, 칸, 12.5, GRAY, MED, first=True)
    tf = 글(s, M + 150, y - 6, 330, 44)
    줄(tf, 값, 21, INK, XBOLD, first=True)
    tf = 글(s, M + 500, y + 4, 292, 40)
    줄(tf, 설명, 11, GRAY, REG, first=True)
    y += 74
각주(s, "2026. 08. 31. 기준 · 건별 내역은 별첨")

# ══════════════════════════════════════════ 2-1. 성과지표 (공식 양식)
s = 세기(장())
머리(s, "성과지표 달성 현황", "2026. 08. 31. 재검수 기준")
칸x = [M, M + 108, M + 386, M + 500, M + 616]
칸w = [104, 274, 110, 110, 176]
머리줄 = ["구분", "성과지표", "목표(A)", "달성(B)", "달성률(B/A)"]
선(s, M, 178, W - M * 2, INK, 1.4)
for x, w, h in zip(칸x, 칸w, 머리줄):
    tf = 글(s, x, 190, w, 22)
    줄(tf, h, 10.5, INK, XBOLD, first=True)
선(s, M, 216, W - M * 2)
지표 = [("정량·필수", "외국인환자 유치", "12 건", "0 건", "0%"),
        ("정량·부가", "사전상담·사후관리", "120 건", "29 건", "24%"),
        ("정량·부가", "환자 만족도", "90 점", "미측정", "10월부터 수집"),
        ("정성", "ICT 사전상담·사후관리 체계", "플랫폼 구축", "구축 완료", "달성"),
        ("정성", "양·한방 협진 연계 치료모델", "협진체계 수립", "수립·운영 중", "달성")]
y = 232
for 구분, 이름, 목표, 달성, 율 in 지표:
    강조 = 이름 in ("외국인환자 유치", "사전상담·사후관리")
    for x, w, v, 굵 in zip(칸x, 칸w, (구분, 이름, 목표, 달성, 율),
                          (REG, MED, REG, XBOLD if 강조 else REG, MED if 강조 else REG)):
        tf = 글(s, x, y, w, 24)
        줄(tf, v, 11 if 굵 != XBOLD else 12.5, INK if 굵 in (MED, XBOLD) else GRAY, 굵, first=True)
    선(s, M, y + 30, W - M * 2)
    y += 44
tf = 글(s, M, y + 14, 760, 26)
줄(tf, "정성지표 둘은 달성했고, 정량은 사전상담만 올라가고 있습니다.", 12, INK, MED, first=True)
각주(s, "세는 단위는 계획서의 「유치 12건 × 10회(사전 5 + 사후 5)」와 같음 · 사전상담 건별 내역과 협진 프로토콜(v1.0)은 별첨")

# ══════════════════════════════════════════ 3. 무엇을 만들었나
s = 세기(장())
머리(s, "플랫폼 구축 완료", "공고문의 6대 ICT 서비스를 모두 구현했습니다. 러시아어·카자흐어를 포함해 6개 언어")
윗줄 = (("07_07.png", "① 병원안내 · 매칭", "제2의료소견을 안내하고 바로 상담 신청으로"),
        ("07_08.png", "② 진료의뢰 · 상담", "어느 병원에서 무슨 치료를 받는지 현지어로"),
        ("07_09.png", "③ 예약상담", "AI 상담 · 메신저 · 문의폼을 한 곳으로"))
아랫줄 = (("11_13.png", "④ 경과 관찰", "환자가 증상을 적으면 AI 가 위험도를 매겨 알림"),
          ("11_15.png", "⑤ 모니터링 · 교육", "암종별 교육자료를 치료단계에 맞춰 자동 배포"),
          ("11_14.png", "⑥ 재이용 예약", "경과와 증상을 보고 재예약이 필요한지 자동 판정"))
tf = 글(s, M, 158, 200, 22)
줄(tf, "사전상담 3종", 11, INK, XBOLD, first=True)
tf = 글(s, M, 328, 200, 22)
줄(tf, "사후관리 3종", 11, INK, XBOLD, first=True)
for 위치, 묶음 in ((182, 윗줄), (352, 아랫줄)):
    for i, (파일, 이름, 설명) in enumerate(묶음):
        x = M + i * 268
        그림(s, 파일, x, 위치, w=248, h=96)
        tf = 글(s, x, 위치 + 102, 248, 34)
        줄(tf, 이름, 10.5, INK, MED, first=True)
        줄(tf, 설명, 9, GRAY, REG, before=2)
각주(s, "healwith.co.kr · 안드로이드 앱 2026.08.13 출시 · iOS 심사 중")

# ══════════════════════════════════════════ 4. 원격협진 실증
s = 세기(장())
머리(s, "원격협진 실증",
     "2026.08.04 카자흐스탄 환자와 화상 상담. 환자·보호자는 현지, 코디네이터·의료진은 한국에서 여섯 화면 동시 접속")
# 두 화면의 가로세로 비가 달라(2.00 / 1.80) 높이를 맞추고 폭을 각자 잡는다
높 = 216
왼폭, 오폭 = int(높 * 2000 / 999), int(높 * 2000 / 1109)   # 432 · 389
사이 = 28
시작 = (W - 왼폭 - 사이 - 오폭) // 2
그림(s, "08_10.png", 시작, 190, w=왼폭, h=높)
그림(s, "08_11.png", 시작 + 왼폭 + 사이, 190, w=오폭, h=높)
tf = 글(s, 시작, 190 + 높 + 14, 왼폭, 24)
줄(tf, "여섯 명이 동시에 접속했습니다", 10.5, GRAY, REG, first=True)
tf = 글(s, 시작 + 왼폭 + 사이, 190 + 높 + 14, 오폭, 24)
줄(tf, "한국어와 러시아어 자막이 말하는 즉시 함께 뜹니다", 10.5, GRAY, REG, first=True)
각주(s, "해외 파트너 협의도 같은 플랫폼에서 진행했습니다 · 07.14 카자흐 파트너 회의 자막 112건")

# ══════════════════════════════════════════ 5. 소견 왕복
s = 세기(장())
머리(s, "제2의료소견 확보·전달", "대학병원 소견을 받아 환자 언어로 옮겨 전달했습니다. 여섯 건 확보 · 여섯 건 전달")
그림(s, "09_12.png", M + 430, 172, w=310)
tf = 글(s, M, 196, 380, 200)
for 제목, 설명, 사이 in (
        ("현지 의무기록을 받아 번역·정리", "러시아어 원본을 영문·한글로", 0),
        ("협진 병원에 의뢰", "면력 · 신촌면력 · 이대서울 · 이대목동", 22),
        ("회신 소견을 환자 언어로 다시", "러시아어 번역본까지 함께 전달", 22)):
    줄(tf, 제목, 13, INK, XBOLD, first=(사이 == 0), before=사이)
    줄(tf, 설명, 10.5, GRAY, REG, before=4)
선(s, M, 400, 380)
tf = 글(s, M, 414, 380, 40)
줄(tf, "환자는 방한 전에 치료 가능 여부와 대략의 비용을 압니다", 11.5, INK, MED, first=True)
각주(s, "오른쪽은 이대서울병원 국제진료센터 회신 (2026.08.14) · 환자 개인정보는 가린 상태")

# ══════════════════════════════════════════ 6. 사전상담 29건
s = 세기(장())
머리(s, "사전상담 29건 · 환자별 내역",
     "「환자 한 건에 대해 한 번 한 일」을 1건으로 셉니다. 계획서의 「유치 12건 × 10회」와 같은 기준입니다")
환자 = [("신장이식", "에티오피아", 5), ("간이식 재수술", "러시아", 9),
        ("자궁경부 고도이형성", "키르기스스탄", 6), ("왼쪽 신장암", "카자흐스탄", 5),
        ("상행결장암", "카자흐스탄", 3), ("저분화 암종", "카자흐스탄", 1)]
최대 = max(n for _, _, n in 환자)
y = 190
for 진단, 나라, n in 환자:
    tf = 글(s, M, y, 200, 24)
    줄(tf, 진단, 11.5, INK, MED, first=True)
    tf = 글(s, M + 208, y + 2, 120, 22)
    줄(tf, 나라, 10, GRAY, REG, first=True)
    막대 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, pt(M + 340), pt(y + 4), pt(n / 최대 * 300), pt(13))
    막대.fill.solid(); 막대.fill.fore_color.rgb = INK
    막대.line.fill.background(); 막대.shadow.inherit = False
    tf = 글(s, M + 652, y + 1, 60, 22)
    줄(tf, f"{n}건", 11, INK, MED, first=True)
    y += 34
선(s, M, y + 6, W - M * 2)
tf = 글(s, M, y + 20, 700, 30)
줄(tf, "접수·검토 / 병원 확인 요청 / 회신 수령 / 영상 상담 / 환자 전달", 10.5, GRAY, REG, first=True)
각주(s, "건별 일자와 증빙 위치는 별첨 「사전상담 건별 내역_4~8월」")

# ══════════════════════════════════════════ 7. 파트너
s = 세기(장())
머리(s, "파트너 네트워크 · 미팅 23회", "카자흐스탄 · 러시아 · 키르기스스탄 · 인도 · 영국 · 튀르키예 · 국내 대학병원")
for i, (파일, 설명) in enumerate((
        ("12_16.png", "06.17  UMIT 국제암센터 (카자흐)"),
        ("12_17.png", "07.06  Madanes 보험사 (러시아)"),
        ("12_18.png", "12_18"))):
    x = M + i * 268
    그림(s, 파일, x, 190, h=146)
    tf = 글(s, x, 344, 248, 22)
    줄(tf, 설명 if i < 2 else "06.19  MedicaTour (러시아)", 10, GRAY, REG, first=True)
선(s, M, 384, W - M * 2)
tf = 글(s, M, 398, 700, 40)
줄(tf, "현지 관행 수수료는 20~30%였습니다. 국내 법정 상한을 설명해 10%로 합의했습니다.", 11.5, INK, MED, first=True)
각주(s, "미팅 23회 = 해외 의료기관·에이전시 21 · 국내 대학병원 2 · 건별 내역은 별첨")

# ══════════════════════════════════════════ 8. 계약 10건
s = 세기(장())
머리(s, "계약 · 협약 · 등록 10건", "환자가 오는 경로를 넓히려고 맺었습니다")
계약 = [("Global Healthcare Opulence", "인도", "04"), ("MedVoyage Global", "영국", "05"),
        ("Medguide.kz", "카자흐스탄", "07"), ("UMIT 국제암센터", "카자흐스탄", "07"),
        ("이화여자대학교 의료원", "대한민국", "07"), ("MedicaTour", "러시아", "07"),
        ("SAULYK Ltd", "영국", "08"), ("신촌세브란스병원", "대한민국", "08"),
        ("Medex Travel", "키르기스스탄", "08"), ("Clinic Navigator", "카자흐스탄", "08")]
for i, (이름, 나라, 달) in enumerate(계약):
    x = M + (i % 2) * 400
    y = 190 + (i // 2) * 46
    tf = 글(s, x, y, 44, 22)
    줄(tf, f"{달}월", 10, GRAY, MED, first=True)
    tf = 글(s, x + 52, y - 2, 220, 24)
    줄(tf, 이름, 12, INK, MED, first=True)
    tf = 글(s, x + 272, y, 110, 22)
    줄(tf, 나라, 10, GRAY, REG, first=True)
    선(s, x, y + 28, 372)
각주(s, "Big Tourism · DeepMed 등도 검토 중")

# ══════════════════════════════════════════ 9. 알리고 있다
s = 세기(장())
머리(s, "현지 홍보 채널", "러시아어·카자흐어 영상과 게시물로 서비스를 알리고 있습니다")
for 파일, 설명, x, 폭, 높이 in (("14_19.png", "홍보영상 · 러시아어", M, 128, 226),
                             ("14_20.png", "홍보영상 · 카자흐어", M + 148, 128, 226),
                             ("14_21.png", "안내 포스터", M + 296, 152, 226),
                             ("14_22.png", "틱톡 @healwith.kz", M + 472, 226, 222)):
    그림(s, 파일, x, 188, w=폭, h=높이)
    tf = 글(s, x, 424, max(폭, 226), 22)
    줄(tf, 설명, 10, GRAY, REG, first=True)
각주(s, "틱톡 영상 8편 누적 조회 1,884회")

# ══════════════════════════════════════════ 10. 안 된 것 (검은 판)
s = 세기(장())
머리(s, "유치 실적 0건", "여기서부터가 오늘 말씀드릴 본론입니다")
선(s, M, 196, W - M * 2, INK, 1.6)
tf = 글(s, M, 226, 240, 110)
줄(tf, "0", 68, INK, HEAVY, first=True)
tf = 글(s, M + 190, 252, 580, 110)
줄(tf, "목표 12건", 12, GRAY, REG, first=True)
줄(tf, "문의는 왔습니다. 여섯 명 모두 의무기록을 보내 왔고", 13.5, INK, REG, before=12)
줄(tf, "병원 검토까지 갔습니다. 그런데 치료 확정이 안 됐습니다.", 13.5, INK, REG, before=5)
선(s, M, 372, W - M * 2)
tf = 글(s, M, 390, 760, 26)
줄(tf, "다음 세 쪽에서 그 사유와 대책을 말씀드립니다.", 11.5, GRAY, REG, first=True)
각주(s, "2026. 08. 31. 기준")

# ══════════════════════════════════════════ 11. 왜 안 됐나
s = 세기(장())
머리(s, "미진 사유", "기능이 모자라서가 아니라, 온 환자의 병세가 이미 깊어서입니다")
이유 = [("문의 대다수가 장기이식·말기암이었습니다",
        "신장이식 · 간이식 재수술 두 건은 공여자가 없어 국내 규정상 수술이 불가능했습니다. "
        "국내 병원에서도 치료가 어렵다는 회신을 받았습니다."),
       ("현지 에이전시가 병원을 직접 보고 보내겠다고 했습니다",
        "카자흐 Medtour 는 담당자 방한 실사 후에 송출을 시작하겠다는 조건을 걸었습니다. 그만큼 초기 송출이 늦어졌습니다."),
       ("의뢰에서 회신까지 걸리는 시간이 중증 환자에게는 깁니다",
        "대학병원 의뢰는 병원마다 접수 절차가 달라 메일과 양식으로 오갑니다. 우리가 그 결과를 다시 플랫폼에 기록합니다.")]
y = 176
for 제목, 설명 in 이유:
    tf = 글(s, M, y, 700, 26)
    줄(tf, 제목, 14, INK, XBOLD, first=True)
    tf = 글(s, M, y + 28, 720, 44)
    줄(tf, 설명, 10.5, GRAY, REG, first=True, spacing=1.35)
    y += 96
각주(s, "진흥원이 2026.07.30 발주한 「비대면진료 기반 사전·사후관리 표준 운영체계」가 세 번째 문제를 겨냥합니다")

# ══════════════════════════════════════════ 12. 어떻게 바꾸나
s = 세기(장())
머리(s, "대책 방향", "환자가 오는 길을 세 갈래로 넓힙니다. 9월부터 실행합니다")
셋 = [("치료 가능한 환자군으로", "종합검진 · 조기암 중심으로 유입 채널을 조정합니다. 말기암 문의는 계속 받되 기대치를 먼저 맞춥니다."),
      ("에이전시를 거치지 않는 길", "현지 검색광고로 환자가 직접 들어오게 합니다. 러시아·CIS 대상 광고는 이미 집행을 시작했습니다."),
      ("보험사 채널", "러시아 Madanes 와 보험 가입자 연계를 열었습니다. 보험 안내 페이지와 전용 계정을 신설했습니다.")]
y = 186
for i, (제목, 설명) in enumerate(셋):
    tf = 글(s, M, y, 60, 40)
    줄(tf, f"0{i+1}", 26, FAINT, HEAVY, first=True)
    tf = 글(s, M + 80, y + 2, 700, 26)
    줄(tf, 제목, 14.5, INK, XBOLD, first=True)
    tf = 글(s, M + 80, y + 30, 700, 40)
    줄(tf, 설명, 10.5, GRAY, REG, first=True, spacing=1.35)
    y += 88
각주(s, "잔여 보조금 약 5,573만원을 현지 광고 · 파트너 초청 · 인건비 · 국외여비에 집중 투입")

# ══════════════════════════════════════════ 13. 남은 석 달
s = 세기(장())
머리(s, "향후 추진일정", "2026. 11. 20. 사업 종료까지 남은 석 달")
달 = [("9월", "첫 유치 확정", "현지 검색광고 본격 집행 · 사전상담 누적 45건"),
      ("10월", "사후관리 가동", "실환자에게 자동 안내 적용 · 만족도 설문 수집 개시 · 유치 누적 8건"),
      ("11월", "목표 달성 · 최종보고", "유치 12건 · 사전상담 120건 · UMIT 통한 현지 경과관찰 데이터 수집")]
y = 196
for i, (때, 제목, 설명) in enumerate(달):
    선(s, M + 30, y + 14, 1.4, INK if i == 0 else FAINT) if False else None
    tf = 글(s, M, y, 70, 30)
    줄(tf, 때, 16, INK, XBOLD, first=True)
    tf = 글(s, M + 100, y, 660, 26)
    줄(tf, 제목, 14, INK, MED, first=True)
    tf = 글(s, M + 100, y + 26, 700, 26)
    줄(tf, 설명, 10.5, GRAY, REG, first=True)
    if i < 2:
        선(s, M, y + 62, W - M * 2)
    y += 86
각주(s, "사업기간 2026. 4. 6. ~ 2026. 11. 20.")

# ══════════════════════════════════════════ 14. 사업비
s = 세기(장())
머리(s, "사업비 집행 현황", "총사업비 대비 23.7% 집행. 국외여비와 위탁사업비가 9월 이후에 몰려 있는 구조입니다")
tf = 글(s, M, 196, 300, 100)
줄(tf, "20,766,610", 30, INK, HEAVY, first=True)
줄(tf, "원 집행액 (총사업비 87,500천원)", 11, GRAY, REG, before=6)
tf = 글(s, M + 400, 196, 400, 100)
줄(tf, "66,733", 30, INK, HEAVY, first=True)
줄(tf, "천원 남음 (사업비 잔액)", 11, GRAY, REG, before=6)
선(s, M, 300, W - M * 2)
쓴것 = [("인건비", "7월분 급여", "5,555,560"), ("홍보비", "온라인 마케팅 위탁 · 영상 2편", "7,000,000"),
       ("운영비", "앱 배포 · 서버 · 데이터베이스 · AI · 서비스 구독료", "507,664"), ("소모품비", "참여기관 소모품 구매 (7월 · 8월)", "1,200,000"),
       ("인건비(현물)", "참여기관 인력 투입 · 자기부담", "6,503,386")]
y = 318
for 비목, 내용, 금액 in 쓴것:
    tf = 글(s, M, y, 110, 22)
    줄(tf, 비목, 11, INK, MED, first=True)
    tf = 글(s, M + 120, y, 420, 22)
    줄(tf, 내용, 10.5, GRAY, REG, first=True)
    tf = 글(s, M + 600, y, 190, 22, align=PP_ALIGN.RIGHT)
    줄(tf, f"{금액} 원", 11, INK, REG, first=True)
    y += 30
각주(s, "2026.08.25 기준 · 08.10 회계법인 제출 「사업비 사용실적보고서」와 같은 기준 · 자체 개발로 외주 개발비 0원")

# ══════════════════════════════════════════ 15. 개발은 지연 없음
s = 세기(장())
머리(s, "개발 추진 현황", "전 항목 계획일정 내 완료. 현재 「운영 및 통제」 단계입니다")
단계 = [("분석", "4월", "현황분석 · GAP분석 · 요구사항 도출"),
       ("설계", "5월", "인터페이스 · 데이터 · 프로그램 · 시험 설계"),
       ("개발", "7월", "6개 언어 · 원격협진 · AI 상담 · 관리자 화면"),
       ("시험", "7월", "통합·시스템 시험 · 자동검사 상시 가동"),
       ("운영", "7월~", "설치·인도 완료 · 운영 및 통제 진행 중")]
y = 194
for 이름, 시점, 내용 in 단계:
    점 = s.shapes.add_shape(MSO_SHAPE.OVAL, pt(M), pt(y + 6), pt(9), pt(9))
    점.fill.solid(); 점.fill.fore_color.rgb = INK
    점.line.fill.background(); 점.shadow.inherit = False
    tf = 글(s, M + 26, y, 80, 24)
    줄(tf, 이름, 12.5, INK, XBOLD, first=True)
    tf = 글(s, M + 116, y + 1, 70, 22)
    줄(tf, 시점, 10.5, GRAY, MED, first=True)
    tf = 글(s, M + 200, y + 1, 560, 22)
    줄(tf, 내용, 10.5, GRAY, REG, first=True)
    y += 46
선(s, M, y + 10, W - M * 2)
tf = 글(s, M, y + 26, 700, 26)
줄(tf, "지연 0주. 계획 일정을 넘긴 항목이 없습니다.", 12.5, INK, MED, first=True)
각주(s, "항목별 계획·추진일정 대조표는 별첨")

# ══════════════════════════════════════════ 16. 마무리 (검은 판)
s = 세기(장(dark=True), dark=True)
tf = 글(s, M, 170, 780, 60)
줄(tf, "만드는 일은 끝났습니다.", 30, ONDARK, XBOLD, first=True)
줄(tf, "남은 것은 환자를 데려오는 일입니다.", 30, ONDARK, XBOLD, before=10)
선(s, M, 306, 96, DIM, 1.6)
tf = 글(s, M, 336, 780, 100)
줄(tf, "플랫폼 구축 완료 · 파트너 23곳 · 계약 10건 · 사전상담 29건", 13, DIM, REG, first=True)
줄(tf, "9월부터 현지 검색광고와 보험사 채널로 유치 전환을 시작합니다.", 13, ONDARK, REG, before=8)

prs.save(OUT)

# ── 테마 손보기: 「새로 치는 글자」는 낱글자 지정이 아니라 테마를 따라간다.
#    빈 발표자료의 기본 테마는 한글이 「맑은 고딕」이라 그대로 두면 대표님이 친 글자가 바뀐다.
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from fix_theme_fonts import 고치기  # noqa: E402

임시 = OUT + ".tmp"
os.replace(OUT, 임시)
결과 = 고치기(임시, OUT, XBOLD, REG)
os.remove(임시)
print(f"저장: {OUT} ({번호[0]}쪽) · 테마 {결과['테마']}개 손봄")

