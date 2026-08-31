"""BeyondK 톤 슬라이드 빌더 — 규격: docs/rules/PPT_STYLE.md

색·판형·글꼴을 한곳에 모아둔다. 새 발표자료는 이 모듈만 import 해서 쓴다.
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색 (원본 PDF 추출값)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF2, 0xF2, 0xF2)
LINE = RGBColor(0xDC, 0xDD, 0xDD)
BODY = RGBColor(0x76, 0x71, 0x71)
MUTED = RGBColor(0x7F, 0x7F, 0x7F)
BRAND = RGBColor(0x0F, 0x76, 0x6E)       # teal-700 · 브랜드 색 (DESIGN.md 영구 고정)
BRAND_PALE = RGBColor(0xCC, 0xFB, 0xF1)  # teal-100 · 짙은 바탕 위 작은 글씨
GREEN_L = RGBColor(0xF0, 0xFD, 0xFA)     # teal-50  · 흐름도 옅은 칸
GREEN_M = RGBColor(0xCC, 0xFB, 0xF1)     # teal-100 · 흐름도 중간 칸
LIME = BRAND                             # 옛 이름. 쓰던 곳이 남아 있어 별칭만 둔다

# ── 글꼴 (무게마다 이름이 따로)
HEAVY = "에스코어 드림 8 Heavy"
XBOLD = "에스코어 드림 7 ExtraBold"
MED = "에스코어 드림 5 Medium"
REG = "에스코어 드림 4 Regular"
LIGHT = "에스코어 드림 3 Light"

W, H = 960, 540
MARGIN = 72


def deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(W), Pt(H)
    return prs


def slide(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BLACK
    return s


def text(s, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h)).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def 글꼴박기(run, 이름):
    """글꼴을 latin·ea(동아시아)·cs 세 칸에 함께 박는다.

    왜: python-pptx 의 font.name 은 <a:latin> «만» 건드린다. 한글은 <a:ea> 를 따라가고
    그게 비면 테마의 script="Hang" 값으로 떨어지는데, 그 기본값이 「맑은 고딕」이다.
    그래서 만들 때는 멀쩡해 보여도 «새로 친 글자»가 맑은 고딕으로 바뀐다
    (2026-08-31 실측 — 중간평가 발표자료에서 실제로 터졌다).

    실측 재현(2026-08-31, python-pptx 1.0.2): 이 함수 «없이» 덱을 뽑으면
    run 7개 / <a:latin> 7개 / **<a:ea> 0개**, 테마의 script="Hang" 은 「맑은 고딕」이었다.
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    rPr = run._r.get_or_add_rPr()
    for 태그 in ("a:ea", "a:cs"):
        el = rPr.find(qn(태그))
        if el is None:
            el = etree.SubElement(rPr, qn(태그))
        el.set("typeface", 이름)


def line(tf, txt, size, color=BLACK, font=REG, first=False, before=0, align=None, spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align:
        p.alignment = align
    p.space_before = Pt(before)
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = txt
    r.font.size, r.font.name = Pt(size), font
    r.font.color.rgb = color
    글꼴박기(r, font)
    return p


def rich(tf, parts, size=12, align=None, before=0, first=False):
    """parts = [(글자, 글꼴, 색), ...] — 한 줄 안에서 강조를 섞을 때."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align:
        p.alignment = align
    p.space_before = Pt(before)
    for txt, font, color in parts:
        r = p.add_run()
        r.text = txt
        r.font.size, r.font.name = Pt(size), font
        r.font.color.rgb = color
        글꼴박기(r, font)
    return p


def text_width(txt, size):
    """형광펜 박스 폭 추정 — 한글·전각은 1em, 그 외는 0.55em."""
    em = sum(1.0 if ord(c) > 0x1100 else 0.55 for c in txt)
    return em * size + 20


def box(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = s.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if radius is not None:
        sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def panel(s, x=MARGIN, y=140, w=W - MARGIN * 2, h=300):
    return box(s, x, y, w, h, PANEL, MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)


def step(s, x, y, w, h, title, lines, fill=PANEL, title_color=BLACK):
    """흐름도 한 칸."""
    box(s, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    dark = (fill == BRAND)          # 짙은 바탕이면 글씨를 뒤집는다
    tf = text(s, x + 10, y + 12, w - 20, h - 20, PP_ALIGN.CENTER)
    line(tf, title, 11, WHITE if dark else title_color, XBOLD, first=True, align=PP_ALIGN.CENTER)
    for t in lines:
        line(tf, t, 9, BRAND_PALE if dark else BODY, REG, before=3, align=PP_ALIGN.CENTER)


def arrow(s, x, y, w=22, h=12):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Pt(x), Pt(y), Pt(w), Pt(h))
    a.fill.solid()
    a.fill.fore_color.rgb = RGBColor(0xA6, 0xA6, 0xA6)
    a.line.fill.background()
    a.shadow.inherit = False
    return a


# ── 슬라이드 5종 ─────────────────────────────────────────

def cover(prs, kicker, title_lines, sub, footer):
    s = slide(prs)
    tf = text(s, 0, 150, W, 200, PP_ALIGN.CENTER)
    line(tf, kicker, 11, MUTED, LIGHT, first=True, align=PP_ALIGN.CENTER)
    for i, t in enumerate(title_lines):
        line(tf, t, 26, BLACK, HEAVY, before=(10 if i == 0 else 2), align=PP_ALIGN.CENTER)
    line(tf, sub, 13, BODY, REG, before=10, align=PP_ALIGN.CENTER)
    tf2 = text(s, 0, 430, W, 40, PP_ALIGN.CENTER)
    line(tf2, footer, 10, MUTED, LIGHT, first=True, align=PP_ALIGN.CENTER)
    return s


def chapter(prs, title):
    s = slide(prs, dark=True)
    tf = text(s, 0, 245, W, 60, PP_ALIGN.CENTER)
    line(tf, title, 28, WHITE, HEAVY, first=True, align=PP_ALIGN.CENTER)
    return s


def statement(prs, title, body_lines, accent=None):
    """가운데 정렬 서술형. body_lines = [(글자, 강조여부), ...]"""
    s = slide(prs)
    tf = text(s, 130, 175, W - 260, 60, PP_ALIGN.CENTER)
    line(tf, title, 24, BLACK, HEAVY, first=True, align=PP_ALIGN.CENTER)
    box(s, W / 2 - 20, 223, 40, 0.8, BLACK)
    tf = text(s, 130, 240, W - 260, 160, PP_ALIGN.CENTER)
    for i, (t, strong) in enumerate(body_lines):
        line(tf, t, 12, BLACK if strong else BODY, XBOLD if strong else REG,
             first=(i == 0), before=(6 if strong else 0), align=PP_ALIGN.CENTER)
    if accent:
        wpx = text_width(accent, 11)
        box(s, W / 2 - wpx / 2, 400, wpx, 24, BRAND)
        tf = text(s, W / 2 - wpx / 2, 400, wpx, 24, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        tf.word_wrap = False
        line(tf, accent, 11, WHITE, XBOLD, first=True, align=PP_ALIGN.CENTER)
    return s


def head(s, kicker, title, subtitle=None):
    tf = text(s, MARGIN + 6, 36, 700, 90)
    line(tf, kicker, 9, MUTED, LIGHT, first=True)
    line(tf, title, 24, BLACK, REG, before=2)
    if subtitle:
        line(tf, subtitle, 11, BLACK, REG, before=4)
    return s


def content(prs, kicker, title, subtitle=None):
    s = slide(prs)
    head(s, kicker, title, subtitle)
    return s


def table(s, rows, x, y, widths, row_h=26, size=10.5, header=True, align_center=False,
          bordered=False):
    """세로선 없는 기본형 표. rows[0] = 머리행."""
    tbl = s.shapes.add_table(len(rows), len(widths), Pt(x), Pt(y),
                             Pt(sum(widths)), Pt(row_h * len(rows))).table
    for i, wv in enumerate(widths):
        tbl.columns[i].width = Pt(wv)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Pt(row_h)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if (header and ri == 0) else WHITE
            cell.margin_left = cell.margin_right = Pt(8)
            cell.margin_top = cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # 줄바꿈(\n)이 있으면 문단이 여러 개 — 전부 서식을 맞춰야 한다
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if align_center else PP_ALIGN.LEFT
                p.line_spacing = 1.15
                for r in p.runs:
                    글꼴 = MED if (header and ri == 0) else REG
                    r.font.size = Pt(size)
                    r.font.name = 글꼴
                    r.font.color.rgb = BLACK
                    글꼴박기(r, 글꼴)
    return tbl


def picture(s, path, x, y, w=None, h=None, border=True):
    """실화면 캡처 삽입. w 또는 h 중 하나만 주면 비율 유지."""
    kw = {}
    if w:
        kw["width"] = Pt(w)
    if h:
        kw["height"] = Pt(h)
    pic = s.shapes.add_picture(path, Pt(x), Pt(y), **kw)
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    return pic


def caption(s, txt, x, y, w, size=9.5, align=PP_ALIGN.LEFT):
    tf = text(s, x, y, w, 16, align)
    line(tf, txt, size, MUTED, LIGHT, first=True, align=align)


def stat(s, x, y, w, value, label, sub=None, accent=False):
    """숫자 카드 — 큰 숫자 + 설명. 텍스트 나열 대신 쓸 것."""
    box(s, x, y, w, 92, BRAND if accent else PANEL)
    ink = WHITE if accent else BLACK
    tf = text(s, x + 14, y + 12, w - 28, 40)
    line(tf, value, 24, ink, HEAVY, first=True)
    line(tf, label, 10.5, ink, MED, before=2)
    if sub:
        line(tf, sub, 9, BRAND_PALE if accent else BODY, LIGHT, before=1)


def band(s, txt, y=468, size=11):
    box(s, MARGIN, y, W - MARGIN * 2, 38, PANEL)
    tf = text(s, MARGIN, y, W - MARGIN * 2, 38, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    line(tf, txt, size, BLACK, REG, first=True, align=PP_ALIGN.CENTER)


def note(s, txt, y=505):
    tf = text(s, MARGIN + 6, y, W - MARGIN * 2, 20)
    line(tf, txt, 8.5, MUTED, LIGHT, first=True)


def highlight(s, txt, x, y, size=11):
    wpx = text_width(txt, size)
    h = size * 2.1
    box(s, x, y, wpx, h, BRAND)
    tf = text(s, x + 6, y, wpx, h, anchor=MSO_ANCHOR.MIDDLE)
    tf.word_wrap = False
    line(tf, txt, size, WHITE, XBOLD, first=True)


def save(prs, 경로, 제목글꼴=None, 본문글꼴=None):
    """저장하면서 «테마의 한글 글꼴»까지 바꾼다. `prs.save()` 대신 이걸 써라.

    왜 낱글자에 박는 것만으론 모자라나: 글꼴박기() 는 «이미 찍은 글자»만 고친다.
    받은 사람이 슬라이드에 **한 글자라도 새로 치면** 그 글자는 테마를 따라가는데,
    빈 발표자료의 테마는 한글이 script="Hang" → 「맑은 고딕」이다.
    (2026-08-31 실측: 뽑은 파일의 theme1.xml 에 `script="Hang" typeface="맑은 고딕"` 2군데.)

    ⚠️ 이 함수는 원래 `fix_theme_fonts` 라는 «저장소에 없는» 모듈을 불렀다. 그대로 가져왔으면
    부르는 순간 ImportError 로 죽는다 — 그래서 밖을 안 부르고 여기서 직접 고치도록 다시 썼다.
    """
    import os
    import shutil
    import zipfile
    from lxml import etree

    prs.save(경로)
    제목 = 제목글꼴 or XBOLD
    본문 = 본문글꼴 or REG
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def 테마고치기(xml_bytes):
        root = etree.fromstring(xml_bytes)
        for 이름, 글꼴 in (("majorFont", 제목), ("minorFont", 본문)):
            블록 = root.find(f".//{{{A}}}fontScheme/{{{A}}}{이름}")
            if 블록 is None:
                continue
            for 태그 in ("latin", "ea", "cs"):
                el = 블록.find(f"{{{A}}}{태그}")
                if el is None:
                    el = etree.SubElement(블록, f"{{{A}}}{태그}")
                el.set("typeface", 글꼴)
            # script="Hang"(한글) 항목도 같이 — 이게 「맑은 고딕」의 진짜 출처다.
            for f in 블록.findall(f"{{{A}}}font"):
                if f.get("script") == "Hang":
                    f.set("typeface", 글꼴)
        return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

    임시 = 경로 + ".tmp"
    os.replace(경로, 임시)
    try:
        with zipfile.ZipFile(임시) as 원본, zipfile.ZipFile(경로, "w", zipfile.ZIP_DEFLATED) as 새것:
            for item in 원본.infolist():
                데이터 = 원본.read(item.filename)
                if item.filename.startswith("ppt/theme/theme") and item.filename.endswith(".xml"):
                    데이터 = 테마고치기(데이터)
                새것.writestr(item, 데이터)
    except Exception:
        # 고치다 실패하면 «원래 파일»이라도 남긴다 — 빈 파일을 남기는 게 제일 나쁘다.
        shutil.move(임시, 경로)
        raise
    finally:
        if os.path.exists(임시):
            os.remove(임시)
    return 경로
