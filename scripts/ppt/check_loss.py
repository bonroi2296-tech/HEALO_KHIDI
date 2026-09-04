# -*- coding: utf-8 -*-
"""저장 뒤 「글이 잘렸나」를 잡는 검사.
바탕판(before)의 문단이 새 판(after)에 그대로 있는지 문단 단위로 대조한다.
제목·띠처럼 «일부러 바꾼 것»은 인자로 받은 목록에 넣어 제외한다.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

def paras(path):
    prs = Presentation(path)
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            if getattr(sh, "has_table", False) and sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        for p in c.text_frame.paragraphs:
                            t = p.text.strip()
                            if t: out.append(t)
            elif getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = p.text.strip()
                    if t: out.append(t)
    return out

a, b = paras(sys.argv[1]), paras(sys.argv[2])
print("문단 수: 바탕판 %d · 새 판 %d" % (len(a), len(b)))
print("글자 수: 바탕판 %d · 새 판 %d" % (sum(len(x) for x in a), sum(len(x) for x in b)))
sa, sb = set(a), set(b)
lost = sorted(sa - sb)
# 잘린 것(앞부분만 남은 것)을 따로 골라낸다 — 가장 위험한 유형
trunc = [x for x in lost if any(y != x and x.startswith(y[:max(8, len(y)//2)]) for y in sb)]
print("\n■ 새 판에서 사라진 문단: %d" % len(lost))
for x in lost: print("   - " + x[:110])
print("\n■ 그중 「앞부분만 남고 잘린 것」으로 의심되는 것: %d" % len(trunc))
for x in trunc: print("   ! " + x[:110])
